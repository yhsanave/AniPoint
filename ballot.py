from show import Show
from pptx import Presentation
import requests
from datetime import date

class Ballot:
    query = '''query ($search: String) {
                    Page(page: 1, perPage: 5) {
                        media(type: ANIME, search: $search) {
                        id
                        title {
                            english
                            romaji
                        }
                        seasonYear
                        coverImage {
                            large
                        }
                    }
                }
            }'''

    mutation = '''mutation ($mediaId: Int, $customLists: [String]) {
                    SaveMediaListEntry (mediaId: $mediaId, status: PLANNING, customLists: $customLists) {
                        id
                    }
                }'''
    
    def __init__(self):
        self.title = "Meeting Title"
        self.date = date.today().strftime("%B %d %Y")
        self.number = "0"
        self.shows = []
        self.prs = None
        self.accessToken = ""
        self.customList = ""
        self.user = ""

    # Generic setters for user defined variables
    def setTitle(self, t: str):
        self.title = t
        return self.title

    def setDate(self, d: str):
        self.date = d
        return self.date

    def setNumber(self, n: str):
        self.number = n
        return self.number

    def setCustomList(self, l: str):
        self.customList = l
        return self.customList

    # Set and Validate Access Token
    def setAccessToken(self, token: str):
        self.accessToken = token
        vQuery = '''query {
            Viewer {
                name
            }
        }
        '''
        try:
            user = requests.post("https://graphql.anilist.co", headers={"Authorization": "Bearer " + self.accessToken}, json={"query": vQuery}).json()["data"]["Viewer"]["name"]
            return f"Logged in as {user}"
        except TypeError:
            if self.accessToken:
                self.accessToken = ""
                return "Invalid Access Token"
            else:
                return "No Account Detected"
    
    # Toggle function for listUpdate
    def toggleListUpdate(self):
        self.listUpdate = not self.listUpdate
        return self.listUpdate

    # Search for anime
    def search(self, term: str):
        result = requests.post("https://graphql.anilist.co", json={"query": Ballot.query, "variables": {"search": term}}).json()["data"]["Page"]["media"]
        return result

    # Generic getters
    def getTitle(self):
        return self.title

    def getDate(self):
        return self.date

    def getNumber(self):
        return self.number

    def getShows(self):
        return self.shows

    # Function for adding a show, creates a new show object and adds it to shows
    def addShow(self, show: Show):
        self.shows.append(show)
        return show

    # Remove the show at a given index in shows
    def removeShow(self, i: Show):
        try:
            return self.shows.pop(self.shows.index(i))
        except ValueError:
            return False

    # Function for reordering the list, takes a show from the given index and inserts it at the new index
    def moveShow(self, index: int, target: int):
        show = self.shows.pop(index)
        return self.shows.insert(target, show)

    # Function for creating the ballot slide, run twice
    def ballotSlide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        ballotLeft = ""
        ballotRight = ""
        for i in range(0,(len(self.shows)//2)):
            ballotLeft += self.shows[i].getTitleENG() + "\n"

        for i in range((len(self.shows)//2), len(self.shows)):
            ballotRight += self.shows[i].getTitleENG() + "\n"

        slide.placeholders[10].text = ballotLeft
        slide.placeholders[11].text = ballotRight

    # Function to create the info slide for each show
    def showSlide(self, show: Show):
        # Check if Subs or Dubs are available
        if show.getSub() and show.getDub():
            subDub = " (Sub/Dub)"
        elif show.getSub():
            subDub = " (Subbed)"
        elif show.getDub():
            subDub = " (Dubbed)"
        else:
            subDub = ""

        # Fill Slide
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[2])
        slide.placeholders[0].text = show.getTitleENG()
        slide.placeholders[10].text = show.getTitleJPN()
        slide.placeholders[11].text = show.getDescription()
        slide.placeholders[12].insert_picture(show.getGifPath())
        slide.placeholders[13].text = f"Episode: {show.getEpisode()}{subDub}\vPremiered: {show.getSeason()}\vSource: {show.getSource()}\vStudio: {show.getStudio()}\vDirector: {show.getDirector()}\vGenres: {show.getGenres()}"

    # Export the ballot to a powerpoint 
    def export(self, templatePath = r".\template.pptx", outputPath = r".\Output.pptx"):
        self.prs = Presentation(templatePath)

        # Title Slide, run only once
        self.prs.slides.add_slide(self.prs.slide_layouts[0])
        self.prs.slides[0].placeholders[0].text = self.title
        self.prs.slides[0].placeholders[1].text = f"{self.date}\vMeeting #{self.number}"

        # Ballot Slide, run once at the start and once at the end
        self.ballotSlide()

        # Make a slide for each show
        for show in self.shows:
            self.showSlide(show)

        # Add second Ballot slide
        self.ballotSlide()

        # Save powerpoint
        self.prs.save(outputPath)

        # If the listUpdate setting is enabled, add all shows to a custom list on anilist
        if self.accessToken:
            for show in self.shows:
                requests.post("https://graphql.anilist.co", headers={"Authorization": "Bearer " + self.accessToken}, json={"query": Ballot.mutation, "variables": {"mediaId": show.getID(), "customLists": [self.customList]}})
